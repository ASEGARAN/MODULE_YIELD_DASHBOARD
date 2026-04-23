#!/usr/bin/env python3
"""
Create PowerPoint presentation for cDPM Recovery Simulation App
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def add_title_slide(prs, title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background
    background = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(7.5))
    background.fill.solid()
    background.fill.fore_color.rgb = RGBColor(26, 35, 126)  # Dark blue
    background.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(200, 200, 200)
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_content_slide(prs, title, bullets, highlight_indices=None):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(26, 35, 126)
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(20)
        p.space_after = Pt(12)

        if highlight_indices and i in highlight_indices:
            p.font.color.rgb = RGBColor(220, 53, 69)  # Red highlight
            p.font.bold = True
        else:
            p.font.color.rgb = RGBColor(50, 50, 50)

    return slide

def add_code_slide(prs, title, code):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(26, 35, 126)
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Code box
    code_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.3), Inches(1.5), Inches(9.4), Inches(5.5))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(40, 44, 52)  # Dark gray
    code_box.line.color.rgb = RGBColor(100, 100, 100)

    # Code text
    code_text = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(9), Inches(5.2))
    tf = code_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = code
    p.font.size = Pt(14)
    p.font.name = "Consolas"
    p.font.color.rgb = RGBColor(200, 200, 200)

    return slide

def add_challenge_slide(prs, title, challenge, solution):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = RGBColor(26, 35, 126)
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    # Challenge box
    chal_label = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(0.5))
    tf = chal_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Challenge:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(220, 53, 69)

    chal_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(2), Inches(9), Inches(1.8))
    chal_box.fill.solid()
    chal_box.fill.fore_color.rgb = RGBColor(255, 240, 240)
    chal_box.line.color.rgb = RGBColor(220, 53, 69)

    chal_text = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(8.6), Inches(1.6))
    tf = chal_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = challenge
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(50, 50, 50)

    # Solution box
    sol_label = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(0.5))
    tf = sol_label.text_frame
    p = tf.paragraphs[0]
    p.text = "Solution:"
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(40, 167, 69)

    sol_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.7), Inches(9), Inches(2.3))
    sol_box.fill.solid()
    sol_box.fill.fore_color.rgb = RGBColor(240, 255, 240)
    sol_box.line.color.rgb = RGBColor(40, 167, 69)

    sol_text = slide.shapes.add_textbox(Inches(0.7), Inches(4.8), Inches(8.6), Inches(2.1))
    tf = sol_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = solution
    p.font.size = Pt(16)
    p.font.color.rgb = RGBColor(50, 50, 50)

    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    add_title_slide(prs,
        "Building & Publishing to MU App Store",
        "cDPM Recovery Simulation Tool\nApril 21, 2026")

    # Slide 2: Overview
    add_content_slide(prs, "What I Built", [
        "cDPM Recovery Simulation Tool for SOCAMM2",
        "Uses Kevin Roos's Hybrid DPM approach",
        "Analyzes three recovery types:",
        "   - New RPx Fix (VERIFIED) - False miscompare detection",
        "   - New BIOS Fix (PROJECTED) - MULTI_BANK_MULTI_DQ timing",
        "   - HW + SOP Fix (PROJECTED) - Hang recovery",
        "Generates interactive HTML reports with heatmaps",
        "Shows single week + 4-week cumulative analysis"
    ])

    # Slide 3: App Features
    add_content_slide(prs, "Key Features", [
        "Hybrid DPM Calculation:",
        "   - MODULE-level (Mod-Sys, Hang): MSNs / Total FIDs × 1M",
        "   - FID-level (DQ, Row, etc.): UFAILs / Total FIDs × 1M",
        "Recovery Simulation with verified & projected fixes",
        "MSN_STATUS × FAILCRAWLER correlation heatmaps",
        "CLI interface with simple arguments",
        "Automatic caching for faster repeated queries",
        "HTML report with Plotly interactive charts"
    ])

    # Slide 4: Publishing Journey
    add_content_slide(prs, "Publishing Journey", [
        "Step 1: Install MU App Store CLI",
        "Step 2: Create pyproject.toml with metadata",
        "Step 3: Structure code as proper Python package",
        "Step 4: Run appstore publish command",
        "Step 5: Test installation from app store",
        "Step 6: Iterate and fix issues",
        "Final: Successfully published v1.0.3!"
    ])

    # Slide 5: Challenge 1 - SSL
    add_challenge_slide(prs, "Challenge 1: SSL Certificate Errors",
        "Error: 'invalid peer certificate: UnknownIssuer'\n\nThe appstore CLI uses 'uv' internally which couldn't connect to GitHub or PyPI due to Micron's self-signed SSL certificates. Every network request failed.",
        "Set environment variables to use system TLS:\n\nexport UV_NATIVE_TLS=1\nexport SSL_CERT_FILE=/etc/pki/ca-trust/extracted/pem/tls-ca-bundle.pem\nexport REQUESTS_CA_BUNDLE=/etc/pki/ca-trust/extracted/pem/tls-ca-bundle.pem\n\nThis tells uv to use the system's native TLS which trusts Micron's CA certificates.")

    # Slide 6: Challenge 2 - Python Version
    add_challenge_slide(prs, "Challenge 2: Python Version Requirements",
        "Error: 'Using Python 3.15 (requires-python: >=3.12)'\n\nThe appstore was trying to download Python 3.15 from GitHub, which failed due to SSL issues. It also ignored our pyproject.toml's requires-python setting.",
        "Force use of system Python with UV environment variables:\n\nexport UV_PYTHON=/path/to/.venv/bin/python\nexport UV_PYTHON_PREFERENCE=only-system\nexport PATH=\"/path/to/.venv/bin:$PATH\"\n\nThis uses the existing Python 3.11 from our virtual environment instead of downloading a new one.")

    # Slide 7: Challenge 3 - Module Not Found
    add_challenge_slide(prs, "Challenge 3: Module Not Included in Package",
        "Error: 'ModuleNotFoundError: No module named cdpm_recovery_sim'\n\nPublishing a single .py file with 'appstore publish script.py' created a wheel, but the actual module code was NOT included in the package!",
        "Create a proper Python package structure:\n\ncdpm_recovery_app/\n├── pyproject.toml\n├── cdpm_recovery_sim.py\n└── cdpm_recovery_sim/        # Package directory\n    ├── __init__.py           # Copy of main script\n    └── __main__.py           # Entry point\n\nThen publish the directory: appstore publish .")

    # Slide 8: Challenge 4 - Entry Point
    add_challenge_slide(prs, "Challenge 4: Wrong Entry Point Auto-Detection",
        "Error: 'Entry points: cdpm-recovery-sim=main:main'\n\nThe appstore auto-detected the wrong entry point, looking for a 'main' module instead of 'cdpm_recovery_sim' module.",
        "Explicitly specify the entry point:\n\nappstore publish . \\\n  --entry-point \"cdpm-recovery-sim=cdpm_recovery_sim:main\"\n\nAlso add [tool.setuptools] to pyproject.toml to exclude unwanted directories:\n\n[tool.setuptools]\npackages = [\"cdpm_recovery_sim\"]")

    # Slide 9: Final Publish Command
    add_code_slide(prs, "Final Working Publish Command",
"""SSL_CERT_FILE=/etc/pki/ca-trust/extracted/pem/tls-ca-bundle.pem \\
REQUESTS_CA_BUNDLE=/etc/pki/ca-trust/extracted/pem/tls-ca-bundle.pem \\
UV_NATIVE_TLS=1 \\
UV_PYTHON=/home/asegaran/MODULE_YIELD_DASHBOARD/.venv/bin/python \\
UV_PYTHON_PREFERENCE=only-system \\
PATH="/home/asegaran/MODULE_YIELD_DASHBOARD/.venv/bin:$PATH" \\
~/.local/bin/appstore publish . \\
  --name cdpm-recovery-sim \\
  --version 1.0.3 \\
  --entry-point "cdpm-recovery-sim=cdpm_recovery_sim:main" \\
  --skip-test \\
  --skip-confirm""")

    # Slide 10: Installation
    add_code_slide(prs, "How Users Install & Run",
"""# Installation (one-time setup)
export UV_NATIVE_TLS=1
export SSL_CERT_FILE=/etc/pki/ca-trust/extracted/pem/tls-ca-bundle.pem
appstore install cdpm-recovery-sim

# Usage
cdpm-recovery-sim --did Y6CP --steps HMB1 --ww 202615

# With multiple design IDs and steps
cdpm-recovery-sim --did Y6CP,Y63N --steps HMB1,QMON --ww 202615

# Custom output file
cdpm-recovery-sim --did Y6CP --steps HMB1 --ww 202615 --output report.html

# Open browser after generation
cdpm-recovery-sim --did Y6CP --steps HMB1 --ww 202615 --browser""")

    # Slide 11: Key Learnings
    add_content_slide(prs, "Key Learnings", [
        "UV_NATIVE_TLS=1 is CRITICAL for Micron network",
        "Single file publishing doesn't work - need package structure",
        "Always specify --entry-point explicitly",
        "Use [tool.setuptools] packages to exclude cache directories",
        "Create a publish.sh script for reproducible deployments",
        "Document SSL requirements in README for users",
        "Test installation on fresh environment before announcing"
    ], highlight_indices=[0, 1, 2])

    # Slide 12: Summary
    add_content_slide(prs, "Summary", [
        "Successfully published cdpm-recovery-sim v1.0.3 to MU App Store",
        "Overcame 4 major challenges with SSL, Python, packaging, and entry points",
        "Created reusable publish.sh script for future updates",
        "Documented installation process in README",
        "Tool is now available for the team to install and use",
        "",
        "Repository: MODULE_YIELD_DASHBOARD/sandbox/cdpm_recovery_app/"
    ])

    # Save
    output_path = "/home/asegaran/MODULE_YIELD_DASHBOARD/sandbox/cdpm_recovery_app/MU_AppStore_Publishing_Journey.pptx"
    prs.save(output_path)
    print(f"Presentation saved: {output_path}")

if __name__ == "__main__":
    main()
