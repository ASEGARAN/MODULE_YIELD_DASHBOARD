#!/usr/bin/env python3
"""Generate VP Readout PPTX with proper waterfall chart visualization."""

import plotly.graph_objects as go
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RgbColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import io

# VP Readout Data (Y6CP SOCAMM2 7500MTPS WW06-15)
data = {
    "total_dpm": 478.0,
    "bios_dpm": 322.1,
    "bios_pct": 67.4,
    "hw_sop_dpm": 19.3,
    "hw_sop_pct": 4.0,
    "recoverable_dpm": 341.4,
    "recoverable_pct": 71.4,
    "remaining_dpm": 136.6,
    "remaining_pct": 28.6,
}

# Create waterfall chart
fig = go.Figure(go.Waterfall(
    name="DPM Recovery",
    orientation="v",
    measure=["absolute", "relative", "relative", "total"],
    x=["Total SLT<br>DPM", "BIOS<br>Recovery", "HW+SOP<br>Recovery", "Customer-Facing<br>DPM"],
    textposition="outside",
    text=[
        f"{data['total_dpm']:.1f}",
        f"-{data['bios_dpm']:.1f}<br>({data['bios_pct']:.0f}%)",
        f"-{data['hw_sop_dpm']:.1f}<br>({data['hw_sop_pct']:.0f}%)",
        f"{data['remaining_dpm']:.1f}"
    ],
    y=[data['total_dpm'], -data['bios_dpm'], -data['hw_sop_dpm'], 0],
    connector={"line": {"color": "rgb(63, 63, 63)", "width": 2}},
    increasing={"marker": {"color": "#FF6B6B"}},  # Red for increases
    decreasing={"marker": {"color": "#4CAF50"}},  # Green for reductions
    totals={"marker": {"color": "#2196F3"}},  # Blue for totals
))

fig.update_layout(
    title={
        'text': "SLT cDPM Recovery Waterfall",
        'y': 0.95,
        'x': 0.5,
        'xanchor': 'center',
        'yanchor': 'top',
        'font': {'size': 24, 'color': '#333', 'family': 'Arial Black'}
    },
    showlegend=False,
    font=dict(size=14, family="Arial"),
    plot_bgcolor='white',
    paper_bgcolor='white',
    yaxis=dict(
        title="DPM",
        showgrid=True,
        gridcolor='lightgray',
        zeroline=True,
        zerolinecolor='gray',
        range=[0, 550]
    ),
    xaxis=dict(
        tickfont=dict(size=12)
    ),
    height=500,
    width=900,
    margin=dict(t=80, b=80, l=80, r=40)
)

# Add annotations for recovery percentages
fig.add_annotation(
    x=0.5, y=-0.15,
    xref='paper', yref='paper',
    text=f"<b>Total Recoverable: {data['recoverable_dpm']:.1f} DPM ({data['recoverable_pct']:.0f}%)</b>  |  <b>Remaining: {data['remaining_dpm']:.1f} DPM ({data['remaining_pct']:.0f}%)</b>",
    showarrow=False,
    font=dict(size=14, color='#333'),
    align='center'
)

# Save waterfall as image
waterfall_bytes = fig.to_image(format="png", scale=2)

# Create PPTX
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Slide 1: Title
slide_layout = prs.slide_layouts[6]  # Blank
slide1 = prs.slides.add_slide(slide_layout)

# Title background shape
title_bg = slide1.shapes.add_shape(
    MSO_SHAPE.RECTANGLE, Inches(0), Inches(2.5), Inches(13.333), Inches(2.5)
)
title_bg.fill.solid()
title_bg.fill.fore_color.rgb = RgbColor(0, 51, 102)
title_bg.line.fill.background()

# Title text
title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.8), Inches(12.333), Inches(1))
title_tf = title_box.text_frame
title_p = title_tf.paragraphs[0]
title_p.text = "Y6CP cDPM Recovery Analysis"
title_p.font.size = Pt(44)
title_p.font.bold = True
title_p.font.color.rgb = RgbColor(255, 255, 255)
title_p.alignment = PP_ALIGN.CENTER

# Subtitle
subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(12.333), Inches(0.6))
subtitle_tf = subtitle_box.text_frame
subtitle_p = subtitle_tf.paragraphs[0]
subtitle_p.text = "SOCAMM2 • 7500MT/s • WW06-15"
subtitle_p.font.size = Pt(28)
subtitle_p.font.color.rgb = RgbColor(200, 200, 200)
subtitle_p.alignment = PP_ALIGN.CENTER

# Slide 2: Waterfall Chart
slide2 = prs.slides.add_slide(slide_layout)

# Slide title
title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
title_tf2 = title_box2.text_frame
title_p2 = title_tf2.paragraphs[0]
title_p2.text = "SLT cDPM Recovery Projection"
title_p2.font.size = Pt(32)
title_p2.font.bold = True
title_p2.font.color.rgb = RgbColor(0, 51, 102)

# Insert waterfall chart image
waterfall_stream = io.BytesIO(waterfall_bytes)
slide2.shapes.add_picture(waterfall_stream, Inches(0.5), Inches(1.0), width=Inches(8))

# Summary metrics box on the right
metrics_box = slide2.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE, Inches(9), Inches(1.2), Inches(4), Inches(5.5)
)
metrics_box.fill.solid()
metrics_box.fill.fore_color.rgb = RgbColor(245, 245, 245)
metrics_box.line.color.rgb = RgbColor(200, 200, 200)

# Metrics title
metrics_title = slide2.shapes.add_textbox(Inches(9.2), Inches(1.4), Inches(3.6), Inches(0.5))
metrics_tf = metrics_title.text_frame
metrics_p = metrics_tf.paragraphs[0]
metrics_p.text = "Recovery Summary"
metrics_p.font.size = Pt(20)
metrics_p.font.bold = True
metrics_p.font.color.rgb = RgbColor(0, 51, 102)
metrics_p.alignment = PP_ALIGN.CENTER

# Metrics content
metrics_content = slide2.shapes.add_textbox(Inches(9.2), Inches(2.0), Inches(3.6), Inches(4.5))
mtf = metrics_content.text_frame
mtf.word_wrap = True

metrics_text = [
    ("Total SLT DPM", f"{data['total_dpm']:.1f}", "#D32F2F"),
    ("", "", ""),
    ("BIOS Recovery", f"{data['bios_dpm']:.1f} ({data['bios_pct']:.0f}%)", "#4CAF50"),
    ("HW+SOP Recovery", f"{data['hw_sop_dpm']:.1f} ({data['hw_sop_pct']:.0f}%)", "#4CAF50"),
    ("", "", ""),
    ("Total Recoverable", f"{data['recoverable_dpm']:.1f} ({data['recoverable_pct']:.0f}%)", "#1976D2"),
    ("", "", ""),
    ("Customer-Facing", f"{data['remaining_dpm']:.1f} ({data['remaining_pct']:.0f}%)", "#FF9800"),
]

for i, (label, value, color) in enumerate(metrics_text):
    if i == 0:
        p = mtf.paragraphs[0]
    else:
        p = mtf.add_paragraph()

    if label:
        p.text = f"{label}: {value}"
        p.font.size = Pt(14)
        p.font.bold = True if label in ["Total SLT DPM", "Total Recoverable", "Customer-Facing"] else False
    else:
        p.text = ""
    p.space_after = Pt(8)

# Legend
legend_box = slide2.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(8), Inches(0.8))
legend_tf = legend_box.text_frame
legend_p = legend_tf.paragraphs[0]
legend_p.text = "BIOS: MULTI_BANK_MULTI_DQ (100%) + Bank/Burst/Periph patterns (50%)  |  HW+SOP: Socket/Debris + SOP compliance"
legend_p.font.size = Pt(11)
legend_p.font.color.rgb = RgbColor(100, 100, 100)

# Slide 3: Step Breakdown Table
slide3 = prs.slides.add_slide(slide_layout)

# Slide title
title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(0.6))
title_tf3 = title_box3.text_frame
title_p3 = title_tf3.paragraphs[0]
title_p3.text = "Step-Level DPM Breakdown"
title_p3.font.size = Pt(32)
title_p3.font.bold = True
title_p3.font.color.rgb = RgbColor(0, 51, 102)

# Table data
table_data = [
    ["Step", "Total DPM", "BIOS", "HW+SOP", "Recoverable", "Remaining"],
    ["HMB1", "287.3", "203.5 (71%)", "11.2 (4%)", "214.7 (75%)", "72.6 (25%)"],
    ["QMON", "190.7", "118.6 (62%)", "8.1 (4%)", "126.7 (66%)", "64.0 (34%)"],
    ["Combined SLT", "478.0", "322.1 (67%)", "19.3 (4%)", "341.4 (71%)", "136.6 (29%)"],
]

# Create table
rows, cols = len(table_data), len(table_data[0])
table = slide3.shapes.add_table(rows, cols, Inches(0.5), Inches(1.2), Inches(12.333), Inches(2.5)).table

# Style table
for i, row in enumerate(table_data):
    for j, cell_text in enumerate(row):
        cell = table.cell(i, j)
        cell.text = cell_text
        cell.text_frame.paragraphs[0].font.size = Pt(14)
        cell.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

        if i == 0:  # Header row
            cell.fill.solid()
            cell.fill.fore_color.rgb = RgbColor(0, 51, 102)
            cell.text_frame.paragraphs[0].font.color.rgb = RgbColor(255, 255, 255)
            cell.text_frame.paragraphs[0].font.bold = True
        elif i == len(table_data) - 1:  # Total row
            cell.fill.solid()
            cell.fill.fore_color.rgb = RgbColor(230, 240, 250)
            cell.text_frame.paragraphs[0].font.bold = True
        else:
            cell.fill.solid()
            cell.fill.fore_color.rgb = RgbColor(255, 255, 255)

# Notes
notes_box = slide3.shapes.add_textbox(Inches(0.5), Inches(4.0), Inches(12), Inches(2.5))
notes_tf = notes_box.text_frame
notes_tf.word_wrap = True

notes = [
    "Key Recovery Mechanisms:",
    "• BIOS: MULTI_BANK_MULTI_DQ failures (100% recovery) + Bank/Burst/Periph patterns (50% recovery)",
    "• HW+SOP: Socket/debris contamination fixes + SOP compliance (HUNG sequence violations)",
    "",
    "Note: 1st pass yield impacted by HW issues (socket and debris contamination). Recovery projections assume fixes deployed."
]

for i, note in enumerate(notes):
    if i == 0:
        p = notes_tf.paragraphs[0]
        p.font.bold = True
    else:
        p = notes_tf.add_paragraph()
    p.text = note
    p.font.size = Pt(12)
    p.font.color.rgb = RgbColor(80, 80, 80)

# Save PPTX
output_path = "/home/asegaran/MODULE_YIELD_DASHBOARD/Y6CP_cDPM_Recovery_VP_Readout_SOCAMM2_7500MT_WW0615.pptx"
prs.save(output_path)
print(f"Saved: {output_path}")
