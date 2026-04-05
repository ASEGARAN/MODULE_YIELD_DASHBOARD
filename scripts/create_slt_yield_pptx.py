"""
Create SOCAMM/SOCAMM2 SLT Yield Issues PowerPoint for VP Debra Bell.
Focus: Y6CP 7500MTPS
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from datetime import datetime
import os
import sys

# Add parent directory to path for imports
sys.path.insert(0, os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from src.fiscal_calendar import get_fiscal_month, get_calendar_year_month


def create_presentation():
    """Create the SLT Yield Issues presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
    prs.slide_height = Inches(7.5)

    # Colors
    MICRON_BLUE = RGBColor(0, 82, 147)  # Micron brand blue
    DARK_GRAY = RGBColor(64, 64, 64)
    RED = RGBColor(220, 53, 69)
    GREEN = RGBColor(40, 167, 69)
    ORANGE = RGBColor(255, 145, 0)

    # Get current workweek info
    today = datetime.now()
    iso_cal = today.isocalendar()
    current_ww = f"{iso_cal[0]}{iso_cal[1]:02d}"
    current_month = get_fiscal_month(current_ww)

    # =========================================================================
    # SLIDE 1: Title + Purpose
    # =========================================================================
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    # Title
    title_box = slide1.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = "SOCAMM / SOCAMM2 SLT Yield Issues"
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = MICRON_BLUE
    title_para.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.333), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = "Current Drivers, Containment, Actions"
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = DARK_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER

    # Scope statement box
    scope_box = slide1.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.333), Inches(1.5))
    scope_frame = scope_box.text_frame
    scope_frame.word_wrap = True

    scope_para1 = scope_frame.paragraphs[0]
    scope_para1.text = "Scope: Y6CP 7500MTPS (SOCAMM2 192GB)"
    scope_para1.font.size = Pt(18)
    scope_para1.font.color.rgb = DARK_GRAY
    scope_para1.alignment = PP_ALIGN.CENTER

    scope_para2 = scope_frame.add_paragraph()
    scope_para2.text = "Definitions: SLT = HMB1 × QMON  |  ELC = HMFN × SLT"
    scope_para2.font.size = Pt(16)
    scope_para2.font.color.rgb = DARK_GRAY
    scope_para2.alignment = PP_ALIGN.CENTER

    # Date
    date_box = slide1.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
    date_frame = date_box.text_frame
    date_para = date_frame.paragraphs[0]
    date_para.text = f"WW{current_ww} ({current_month}) | Prepared for: Debra Bell, VP"
    date_para.font.size = Pt(14)
    date_para.font.color.rgb = DARK_GRAY
    date_para.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 2: Executive Snapshot
    # =========================================================================
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])

    # Title
    add_slide_title(slide2, "Executive Snapshot", MICRON_BLUE)

    # Three key bullets
    content_box = slide2.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Bullet 1: What's impacted
    para1 = content_frame.paragraphs[0]
    para1.text = "Current Status"
    para1.font.size = Pt(24)
    para1.font.bold = True
    para1.font.color.rgb = RED
    para1.space_after = Pt(6)

    detail1 = content_frame.add_paragraph()
    detail1.text = "Y6CP 7500MTPS SOCAMM2: SLT yield at 97.34% — meeting short-term target (97%) but significant gap to long-term target of 99.5% (-2.16%). Continued focus required."
    detail1.font.size = Pt(18)
    detail1.font.color.rgb = DARK_GRAY
    detail1.level = 1
    detail1.space_after = Pt(24)

    # Bullet 2: Top drivers
    para2 = content_frame.add_paragraph()
    para2.text = "Top Drivers"
    para2.font.size = Pt(24)
    para2.font.bold = True
    para2.font.color.rgb = ORANGE
    para2.space_after = Pt(6)

    detail2 = content_frame.add_paragraph()
    detail2.text = "• Hung / CPU Thermal issues\n• Mod-Sys failures (decode limitations)\n• Blank pre-test failures"
    detail2.font.size = Pt(18)
    detail2.font.color.rgb = DARK_GRAY
    detail2.level = 1
    detail2.space_after = Pt(24)

    # Bullet 3: What's improved
    para3 = content_frame.add_paragraph()
    para3.text = "What's Improved"
    para3.font.size = Pt(24)
    para3.font.bold = True
    para3.font.color.rgb = GREEN
    para3.space_after = Pt(6)

    detail3 = content_frame.add_paragraph()
    detail3.text = "New BIOS implementation enabling better decode and reducing Mod-Sys failures. Retest shows high recoverability on Mod-Sys classified units."
    detail3.font.size = Pt(18)
    detail3.font.color.rgb = DARK_GRAY
    detail3.level = 1

    # =========================================================================
    # SLIDE 3: Quantify the Gap
    # =========================================================================
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide3, "Quantify the Gap: Y6CP 7500MTPS Yield vs Target", MICRON_BLUE)

    # Data table placeholder
    content_box = slide3.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Yield targets section
    para1 = content_frame.paragraphs[0]
    para1.text = "Target Yields (C2 Y6CP 7.5Gbps)"
    para1.font.size = Pt(20)
    para1.font.bold = True
    para1.font.color.rgb = MICRON_BLUE
    para1.space_after = Pt(12)

    targets = content_frame.add_paragraph()
    targets.text = """• HMFN Target: 99.00%
• SLT Short-term Target: 97.00% (Mar'26+)
• SLT Long-term Target: 99.50%
• ELC Target: 96.03% (Mar+)"""
    targets.font.size = Pt(16)
    targets.font.color.rgb = DARK_GRAY
    targets.space_after = Pt(24)

    # Current performance
    para2 = content_frame.add_paragraph()
    para2.text = "Current Performance (WW202605-WW202614)"
    para2.font.size = Pt(20)
    para2.font.bold = True
    para2.font.color.rgb = MICRON_BLUE
    para2.space_after = Pt(12)

    perf = content_frame.add_paragraph()
    perf.text = """• HMFN FY: 98.42% (vs 99.00% target) — Gap: -0.58%
• HMB1 FY: 98.42%
• QMON FY: 98.90%
• SLT FY: 97.34% (vs 97.00% target) — Gap: +0.34%
• ELC FY: 95.80% (vs 96.03% target) — Gap: -0.23%

Volume: 30,164 units tested at HMFN"""
    perf.font.size = Pt(16)
    perf.font.color.rgb = DARK_GRAY
    perf.space_after = Pt(24)

    # Key observation
    obs = content_frame.add_paragraph()
    obs.text = "Gap to Long-term Target"
    obs.font.size = Pt(20)
    obs.font.bold = True
    obs.font.color.rgb = RED
    obs.space_after = Pt(8)

    obs_detail = content_frame.add_paragraph()
    obs_detail.text = """• SLT at 97.34% — Gap to long-term 99.5% target: -2.16%
• Top SLT loss drivers: Hung/CPU Thermal, Mod-Sys, Blank pre-test
• HMFN yield gap (-0.58%) also contributing to ELC miss
• WW202612 anomaly: HMFN dropped to 65.04% (under investigation)"""
    obs_detail.font.size = Pt(14)
    obs_detail.font.color.rgb = DARK_GRAY

    # Source note
    note = content_frame.add_paragraph()
    note.text = "\nSource: Module Yield Dashboard - Data as of WW202614"
    note.font.size = Pt(12)
    note.font.italic = True
    note.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 4: Primary Driver 1 - Hung / CPU Thermal
    # =========================================================================
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide4, "Primary Driver 1: Hung / CPU Thermal", MICRON_BLUE)

    content_box = slide4.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Root cause
    para1 = content_frame.paragraphs[0]
    para1.text = "Root Cause Analysis"
    para1.font.size = Pt(20)
    para1.font.bold = True
    para1.font.color.rgb = RED
    para1.space_after = Pt(12)

    rc = content_frame.add_paragraph()
    rc.text = """• CPU temperature remains high (~90°C) on short engineering test program
• Root cause leaning toward equipment-related issue, not test program dependent
• Boards under FA lockdown due to elevated CPU temperatures"""
    rc.font.size = Pt(18)
    rc.font.color.rgb = DARK_GRAY
    rc.space_after = Pt(24)

    # Investigation
    para2 = content_frame.add_paragraph()
    para2.text = "Ongoing Investigation"
    para2.font.size = Pt(20)
    para2.font.bold = True
    para2.font.color.rgb = ORANGE
    para2.space_after = Pt(12)

    inv = content_frame.add_paragraph()
    inv.text = """• Data collection via IPMI monitoring in progress
• Equipment forum investigation ongoing
• Thermal solution evaluation underway"""
    inv.font.size = Pt(18)
    inv.font.color.rgb = DARK_GRAY
    inv.space_after = Pt(24)

    # Impact
    para3 = content_frame.add_paragraph()
    para3.text = "Yield Impact"
    para3.font.size = Pt(20)
    para3.font.bold = True
    para3.font.color.rgb = MICRON_BLUE
    para3.space_after = Pt(12)

    impact = content_frame.add_paragraph()
    impact.text = "[INSERT FAILCRAWLER DATA: Hung failure contribution to SLT fallout]"
    impact.font.size = Pt(16)
    impact.font.italic = True
    impact.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 5: Primary Driver 2 - Mod-Sys + Decode
    # =========================================================================
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide5, "Primary Driver 2: Mod-Sys + Decode Confidence", MICRON_BLUE)

    content_box = slide5.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(11.833), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Challenge
    para1 = content_frame.paragraphs[0]
    para1.text = "Key Challenge: Decode Limitation"
    para1.font.size = Pt(20)
    para1.font.bold = True
    para1.font.color.rgb = RED
    para1.space_after = Pt(12)

    challenge = content_frame.add_paragraph()
    challenge.text = """• Mod-Sys failures indicate system-level issues with limited decode capability
• Current BIOS encryption limits failure analysis depth
• System-to-DQ mapping needed for accurate failure attribution"""
    challenge.font.size = Pt(18)
    challenge.font.color.rgb = DARK_GRAY
    challenge.space_after = Pt(24)

    # Improvement
    para2 = content_frame.add_paragraph()
    para2.text = "BIOS Improvement"
    para2.font.size = Pt(20)
    para2.font.bold = True
    para2.font.color.rgb = GREEN
    para2.space_after = Pt(12)

    improvement = content_frame.add_paragraph()
    improvement.text = """• New BIOS implementation enabled to ensure failures are decoded correctly
• Encryption-disabled BIOS enables deeper debugging
• Improved system-to-DQ mapping for accurate failure classification"""
    improvement.font.size = Pt(18)
    improvement.font.color.rgb = DARK_GRAY
    improvement.space_after = Pt(24)

    # Retest results
    para3 = content_frame.add_paragraph()
    para3.text = "Retest Results"
    para3.font.size = Pt(20)
    para3.font.bold = True
    para3.font.color.rgb = MICRON_BLUE
    para3.space_after = Pt(12)

    retest = content_frame.add_paragraph()
    retest.text = """• Retest shows high recoverability on Mod-Sys classified units
• Ongoing deep dive to quantify true DRAM vs system failures
• [INSERT: Mod-Sys retest pass rate from FAILCRAWLER analysis]"""
    retest.font.size = Pt(18)
    retest.font.color.rgb = DARK_GRAY

    # =========================================================================
    # SLIDE 6: Containment, Owners, and Exec Asks
    # =========================================================================
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])

    add_slide_title(slide6, "Containment, Owners & Executive Asks", MICRON_BLUE)

    content_box = slide6.shapes.add_textbox(Inches(0.75), Inches(1.8), Inches(5.5), Inches(5))
    content_frame = content_box.text_frame
    content_frame.word_wrap = True

    # Containment actions
    para1 = content_frame.paragraphs[0]
    para1.text = "Containment Actions"
    para1.font.size = Pt(18)
    para1.font.bold = True
    para1.font.color.rgb = MICRON_BLUE
    para1.space_after = Pt(8)

    actions = content_frame.add_paragraph()
    actions.text = """• Taskforce ongoing to address yield impact
• Pre-testing issue investigation active
• New BIOS deployment in progress
• Thermal monitoring enhanced"""
    actions.font.size = Pt(14)
    actions.font.color.rgb = DARK_GRAY
    actions.space_after = Pt(16)

    # Owners
    para2 = content_frame.add_paragraph()
    para2.text = "Action Owners"
    para2.font.size = Pt(18)
    para2.font.bold = True
    para2.font.color.rgb = MICRON_BLUE
    para2.space_after = Pt(8)

    owners = content_frame.add_paragraph()
    owners.text = """• Hung/Thermal: [Equipment Team]
• Mod-Sys/Decode: [BIOS/Debug Team]
• Blank Pre-test: [Test Engineering]
• Tracking/DOE: [Yield Engineering]"""
    owners.font.size = Pt(14)
    owners.font.color.rgb = DARK_GRAY

    # Right column - Exec Asks
    ask_box = slide6.shapes.add_textbox(Inches(6.5), Inches(1.8), Inches(6), Inches(5))
    ask_frame = ask_box.text_frame
    ask_frame.word_wrap = True

    para3 = ask_frame.paragraphs[0]
    para3.text = "Executive Asks"
    para3.font.size = Pt(18)
    para3.font.bold = True
    para3.font.color.rgb = RED
    para3.space_after = Pt(12)

    asks = ask_frame.add_paragraph()
    asks.text = """1. Equipment team resource prioritization
   for thermal root cause

2. Expedited BIOS release approval for
   decode improvement

3. Cross-functional alignment on
   Mod-Sys classification criteria"""
    asks.font.size = Pt(14)
    asks.font.color.rgb = DARK_GRAY
    asks.space_after = Pt(24)

    # Timeline
    para4 = ask_frame.add_paragraph()
    para4.text = "Target Timeline"
    para4.font.size = Pt(18)
    para4.font.bold = True
    para4.font.color.rgb = GREEN
    para4.space_after = Pt(8)

    timeline = ask_frame.add_paragraph()
    timeline.text = """• Thermal RC: +2 weeks
• New BIOS deployment: +1 week
• Yield target achievement: [TBD based on RC]"""
    timeline.font.size = Pt(14)
    timeline.font.color.rgb = DARK_GRAY

    # Save presentation
    output_path = "/home/asegaran/MODULE_YIELD_DASHBOARD/output/SOCAMM_SLT_Yield_Issues_VP_Deck.pptx"
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs.save(output_path)
    print(f"Presentation saved to: {output_path}")
    return output_path


def add_slide_title(slide, title_text, color):
    """Add a consistent title to a slide."""
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(12.333), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.size = Pt(32)
    title_para.font.bold = True
    title_para.font.color.rgb = color

    # Add underline
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.5), Inches(1.4), Inches(12.333), Inches(0.03)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()


if __name__ == "__main__":
    output = create_presentation()
    print(f"\nPresentation created successfully!")
    print(f"Location: {output}")
    print("\nSlide structure:")
    print("  1. Title + Purpose (30 sec)")
    print("  2. Executive Snapshot (1 min)")
    print("  3. Quantify the Gap (1-1.5 min)")
    print("  4. Primary Driver 1: Hung/CPU Thermal (1.5 min)")
    print("  5. Primary Driver 2: Mod-Sys + Decode (1.5 min)")
    print("  6. Containment, Owners & Exec Asks (2 min)")
    print("\nNote: Update placeholders with actual dashboard data before presenting.")
